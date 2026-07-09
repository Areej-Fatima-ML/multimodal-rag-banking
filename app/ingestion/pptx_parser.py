from pptx import Presentation
import io
from PIL import Image


def extract_text_from_pptx(file_path: str) -> list:
    """
    Extract all text from PowerPoint slides
    Combines all text shapes from each slide
    Returns list of text dicts with content and metadata
    """
    texts = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides):
        slide_text = ""

        # Extract text from all shapes in slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    slide_text += shape.text + "\n"

        # Only add non-empty slides
        if slide_text.strip():
            texts.append({
                "content": slide_text,
                "slide": slide_num + 1,
                "type": "text",
                "source": file_path
            })

    return texts


def extract_tables_from_pptx(file_path: str) -> list:
    """
    Extract all tables from PowerPoint slides
    Converts table data to readable text format
    Returns list of table dicts with content and metadata
    """
    tables = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_table:
                table_text = ""

                for row in shape.table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_text += " | ".join(row_data) + "\n"

                # Only add non-empty tables
                if table_text.strip():
                    tables.append({
                        "content": table_text,
                        "slide": slide_num + 1,
                        "type": "table",
                        "source": file_path
                    })

    return tables


def extract_images_from_pptx(file_path: str) -> list:
    """
    Extract all images from PowerPoint slides
    Returns list of image dicts with PIL image and metadata
    """
    images = []
    prs = Presentation(file_path)

    for slide_num, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Shape type 13 = Picture
            if shape.shape_type == 13:
                image_data = shape.image.blob
                image = Image.open(io.BytesIO(image_data))

                
                if image.width >= 10 and image.height >= 10:
                    images.append({
                        "image": image,
                        "slide": slide_num + 1,
                        "type": "image",
                        "source": file_path
                    })

    return images


def parse_pptx(file_path: str) -> dict:
    """
    Parse complete PowerPoint file
    Extracts text, tables and images from all slides
    Returns dict with all extracted content and metadata
    """
    texts = extract_text_from_pptx(file_path)
    tables = extract_tables_from_pptx(file_path)
    images = extract_images_from_pptx(file_path)

    return {
        "file_path": file_path,
        "texts": texts,
        "tables": tables,
        "images": images,
        "total_slides": len(Presentation(file_path).slides),
        "total_images": len(images)
    }